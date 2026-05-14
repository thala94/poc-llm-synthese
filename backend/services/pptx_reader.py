from pptx import Presentation

def read_pptx(file_path):
    presentation = Presentation(file_path)
    text = ""

    for i, slide in enumerate(presentation.slides, start=1):
        text += f"\n--- Diapositive {i} ---\n"

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text += shape.text + "\n"

    return text